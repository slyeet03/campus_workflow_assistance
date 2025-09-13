import os
import json
import requests
from flask import Flask, jsonify, request, make_response
from pptx import Presentation
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# -------------------- CONFIG --------------------
UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "AIzaSyCozpuNd_caNfR0Xo-htgW5J6RpqF9yD1c")
GEMINI_URL = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.0-flash-lite:generateContent?key={GOOGLE_API_KEY}"

app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# CORS headers - SIMPLE FIX
@app.after_request
def after_request(response):
    response.headers['Access-Control-Allow-Origin'] = '*'
    response.headers['Access-Control-Allow-Headers'] = 'Content-Type,Authorization'
    response.headers['Access-Control-Allow-Methods'] = 'GET,PUT,POST,DELETE,OPTIONS'
    return response

# Handle preflight requests
@app.before_request
def handle_preflight():
    if request.method == "OPTIONS":
        response = make_response()
        response.headers['Access-Control-Allow-Origin'] = '*'
        response.headers['Access-Control-Allow-Headers'] = 'Content-Type,Authorization'
        response.headers['Access-Control-Allow-Methods'] = 'GET,PUT,POST,DELETE,OPTIONS'
        return response

# -------------------- IN-MEMORY STORAGE --------------------
execution_log = []
timetable = []
pending_quizzes = {}  # quiz_id -> quiz questions waiting for teacher edit
published_quizzes = {}  # quiz_id -> final published quiz

# -------------------- UTILITIES --------------------
def log_action(action):
    execution_log.append(action)

def ppt_to_text(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def call_gemini_api(prompt, max_output_tokens=500):
    """Call Gemini API with correct format"""
    headers = {"Content-Type": "application/json"}
    data = {
        "contents": [{
            "parts": [{"text": prompt}]
        }],
        "generationConfig": {
            "temperature": 0.2,
            "maxOutputTokens": max_output_tokens
        }
    }
    try:
        response = requests.post(GEMINI_URL, headers=headers, json=data)
        print(f"Gemini API status code: {response.status_code}")
        
        if response.status_code == 200:
            result = response.json()
            print(f"Gemini API response: {result}")
            # Extract text from response
            text = result.get("candidates", [{}])[0].get("content", {}).get("parts", [{}])[0].get("text", "")
            if text:
                return text
            else:
                print("No text found in Gemini response")
                return "Error: No content generated"
        else:
            print("Gemini API error:", response.text)
            return "Error: could not generate content"
    except Exception as e:
        print(f"Exception calling Gemini API: {e}")
        return "Error: API call failed"

def summarize_text(text):
    prompt = f"Summarize the following lecture into 5-7 bullet points:\n{text}"
    return call_gemini_api(prompt)

def generate_quiz(summary_text, n_questions, marks_per_question):
    prompt = f"""
    Create a quiz from this lecture summary:
    {summary_text}

    Number of questions: {n_questions}
    Marks per question: {marks_per_question}

    Return JSON with keys: question, type (short_answer/mcq), marks, options (if MCQ)
    """
    print(f"Calling Gemini API with prompt: {prompt[:100]}...")
    response_text = call_gemini_api(prompt)
    print(f"Gemini response: {response_text[:200]}...")
    
    try:
        quiz = json.loads(response_text)
        print(f"Successfully parsed quiz: {quiz}")
    except Exception as e:
        print(f"Failed to parse JSON: {e}")
        print(f"Raw response: {response_text}")
        # Create a better fallback quiz
        quiz = []
        for i in range(n_questions):
            quiz.append({
                "question": f"Question {i+1}: Based on the lecture content, explain the key concept of {summary_text[:50]}...",
                "type": "short_answer",
                "marks": marks_per_question,
                "options": []
            })
    return quiz

# -------------------- ROUTES --------------------
@app.route("/test", methods=["GET"])
def test():
    return jsonify({"message": "Backend is working!", "status": "success"})
@app.route("/upload_ppt", methods=["POST"])
def upload_ppt():
    try:
        if 'ppt' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400
        
        file = request.files["ppt"]
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        n_questions = int(request.form.get("n_questions", 5))
        marks = int(request.form.get("marks", 1))

        # Ensure uploads directory exists
        os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
        
        file_path = os.path.join(app.config["UPLOAD_FOLDER"], file.filename)
        file.save(file_path)
        
        print(f"File saved to: {file_path}")  # Debug print

        text = ppt_to_text(file_path)
        summary = summarize_text(text)
        quiz = generate_quiz(summary, n_questions, marks)

        quiz_id = file.filename
        pending_quizzes[quiz_id] = quiz

        log_action(f"Generated quiz for {file.filename} (pending teacher review)")
        return jsonify({"quiz_id": quiz_id, "summary": summary, "quiz": quiz})
    
    except Exception as e:
        print(f"Error in upload_ppt: {str(e)}")  # Debug print
        return jsonify({"error": str(e)}), 500

@app.route("/edit_quiz", methods=["POST"])
def edit_quiz():
    data = request.json
    quiz_id = data["quiz_id"]
    edits = data.get("edits", [])

    if quiz_id not in pending_quizzes:
        return jsonify({"error": "Quiz not found"}), 404

    quiz = pending_quizzes[quiz_id]
    for edit in edits:
        idx = edit["index"]
        new_q = edit["new_question"]
        if 0 <= idx < len(quiz):
            quiz[idx]["question"] = new_q

    pending_quizzes[quiz_id] = quiz
    log_action(f"Teacher edited quiz {quiz_id}")
    return jsonify({"quiz_id": quiz_id, "quiz": quiz})

@app.route("/publish_quiz", methods=["POST"])
def publish_quiz():
    data = request.json
    quiz_id = data["quiz_id"]

    if quiz_id not in pending_quizzes:
        return jsonify({"error": "Quiz not found"}), 404

    published_quizzes[quiz_id] = pending_quizzes.pop(quiz_id)
    log_action(f"Published quiz {quiz_id}")
    return jsonify({"quiz_id": quiz_id, "message": f"Quiz {quiz_id} published successfully!"})

@app.route("/schedule_class", methods=["POST"])
def schedule_class():
    data = request.json
    year = data["year"]
    branch = data["branch"]
    section = data["section"]
    day = data.get("day")
    mode = data.get("mode", "offline")

    hour = find_free_slot(year, branch, section, day)
    if hour is None:
        return jsonify({"error": "No free slot available"}), 400

    entry = {"year": year, "branch": branch, "section": section, "hour": hour, "mode": mode}
    timetable.append(entry)

    meeting_link = ""
    if mode.lower() == "online":
        meeting_link = create_fake_meeting_link()
        entry["meeting_link"] = meeting_link

    log_action(f"Scheduled class for {year}-{branch}-{section} at {hour}:00 ({mode})")
    return jsonify({"slot": hour, "meeting_link": meeting_link})

@app.route("/execution_log", methods=["GET"])
def get_log():
    return jsonify(execution_log)

@app.route("/timetable", methods=["GET"])
def get_timetable():
    return jsonify(timetable)

@app.route("/pending_quizzes", methods=["GET"])
def get_pending_quizzes():
    return jsonify(pending_quizzes)

@app.route("/published_quizzes", methods=["GET"])
def get_published_quizzes():
    return jsonify(published_quizzes)

# -------------------- SLOT UTILITIES --------------------
def find_free_slot(year, branch, section, day=None):
    for hour in range(9, 18):
        slot_taken = any(
            c["year"] == year and c["branch"] == branch and c["section"] == section and c["hour"] == hour
            for c in timetable
        )
        if not slot_taken:
            return hour
    return None

def create_fake_meeting_link():
    return "https://meet.google.com/link-demo"

# -------------------- RUN --------------------
if __name__ == "__main__":
    port = int(os.getenv("PORT", 5001))
    debug = os.getenv("FLASK_ENV") != "production"
    print(f"Starting Flask app on port {port}...")
    app.run(debug=debug, host='0.0.0.0', port=port)
